import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
import webcolors
import os
import pandas as pd
import openpyxl
from colorama import Fore, Style


# Function to create a new presentation
def create_presentation():
    prs = Presentation()
    return prs


# Function to add a new slide
def add_slide(prs, slide_layout):
    slide = prs.slides.add_slide(slide_layout)
    return slide


# Function to add a title slide
def add_title_slide(slide, title_text):
    title = slide.shapes.title
    title.text = title_text

    # Remove subtitle placeholder
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ''


# Function to add a title and content slide
def add_title_content_slide(slide):
    title = slide.shapes.title
    title.text = input(Fore.YELLOW + Style.BRIGHT + "Input title: ")

    while True:
        content_choice = input(Fore.YELLOW + Style.BRIGHT + 'Do you want to insert a graph or a table? (graph/table): ')

        if content_choice == 'graph':
            print()
            print(Fore.LIGHTYELLOW_EX + Style.BRIGHT + 'Graph')
            while True:
                chart_data = CategoryChartData()
                chart_data.categories = input(
                    Fore.MAGENTA + Style.BRIGHT + 'Enter graph categories (comma-separated): ').split(',')
                series_name = input(Fore.MAGENTA + Style.BRIGHT + 'Enter series name for the graph: ')
                series_values = [float(value) for value in
                                 input(Fore.MAGENTA + Style.BRIGHT + 'Enter series values (comma-separated): ').split(
                                     ',')]
                if len(chart_data.categories) == len(series_values):
                    chart_data.add_series(series_name, series_values)
                    break
                else:
                    print(Fore.RED + Style.BRIGHT + 'Number of categories and values must be the same.')

            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
            chart.has_legend = True
            chart.legend.position = 2  # bottom
            chart.legend.include_in_layout = False

            # Format the pie chart
            for i, point in enumerate(chart.series[0].points):
                while True:
                    color_code = input(Fore.MAGENTA + Style.BRIGHT + f"Enter color for Pie Chart:")
                    try:
                        rgb_color = webcolors.hex_to_rgb(color_code)
                        fill = point.format.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(rgb_color.red, rgb_color.green, rgb_color.blue)
                        break
                    except ValueError:
                        try:
                            rgb_color = webcolors.name_to_rgb(color_code)
                            fill = point.format.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(rgb_color.red, rgb_color.green, rgb_color.blue)
                            break
                        except ValueError:
                            print(Fore.RED + Style.BRIGHT + f'Invalid color code: {color_code}')
            break  # Exit the loop after processing the graph input

        elif content_choice == 'table':
            print()
            print(Fore.LIGHTYELLOW_EX + Style.BRIGHT + 'Table')
            table_data_source = input(
                Fore.MAGENTA + Style.BRIGHT + 'Enter the source of the table data (input/excel): ')
            if table_data_source == 'input':
                num_rows = get_positive_integer_input(Fore.CYAN + Style.BRIGHT + 'Enter the number of rows in the table: ')
                num_cols = get_positive_integer_input(Fore.CYAN + Style.BRIGHT + 'Enter the number of columns in the table: ')
                table = slide.shapes.add_table(num_rows, num_cols, Inches(3), Inches(2), Inches(4), Inches(3)).table
                for row in range(num_rows):
                    print()
                    for col in range(num_cols):
                        while True:
                            cell_value = input(
                                Fore.MAGENTA + Style.BRIGHT + f'Enter value for cell ({row + 1}, {col + 1}): ')
                            if cell_value.strip():
                                table.cell(row, col).text = cell_value
                                break
                            else:
                                print(Fore.RED + Style.BRIGHT + 'Cell value cannot be empty.')
            elif table_data_source == 'excel':
                file_path = input(Fore.CYAN + Style.BRIGHT + 'Enter the path of the Excel file: ')
                sheet_name = input(Fore.CYAN + Style.BRIGHT + 'Enter the name of the Excel sheet: ')

                try:
                    workbook = openpyxl.load_workbook(file_path, read_only=True)
                    sheet = workbook[sheet_name]
                    num_rows = sheet.max_row
                    num_cols = sheet.max_column
                    table = slide.shapes.add_table(num_rows, num_cols, Inches(3), Inches(2), Inches(4),
                                                   Inches(3)).table
                    for row in range(num_rows):
                        for col in range(num_cols):
                            cell_value = sheet.cell(row + 1, col + 1).value
                            if cell_value is not None:
                                table.cell(row, col).text = str(cell_value)
                except FileNotFoundError:
                    print(Fore.RED + Style.BRIGHT + f'File not found at path: {file_path}')
                except KeyError:
                    print(Fore.RED + Style.BRIGHT + f'No sheet found with name: {sheet_name}')
            break  # Exit the loop after processing the table input
        else:
            print(Fore.RED + Style.BRIGHT + 'Invalid choice! Please enter either "graph" or "table".')


# Function to add a section header slide with optional image
def add_section_header_slide(slide):
    title = slide.shapes.title
    title.text = input(Fore.YELLOW + Style.BRIGHT + 'Enter title: ')

    # Remove subtitle placeholder
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ''

    while True:
        image = input(Fore.MAGENTA + Style.BRIGHT + 'Do you want to insert an image? (y/n): ')
        if image == 'y':
            while True:
                img_path = input(Fore.MAGENTA + Style.BRIGHT + 'Enter image path: ')
                try:
                    pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1), height=Inches(3),width=Inches(3))
                    break  # Valid image path, exit the loop
                except FileNotFoundError:
                    print(Fore.RED + Style.BRIGHT + f'Error: Image file not found at "{img_path}"')
        elif image == 'n':
            break  # No image selected, exit the loop
        else:
            print(Fore.RED + Style.BRIGHT + 'Invalid choice! Please enter either "y" or "n".')


# Function to determine text color based on background color brightness
def determine_text_color(bg_color):
    r, g, b = bg_color
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    if brightness > 125:
        return '000000'  # Black
    else:
        return 'FFFFFF'  # White


# Function to get a positive integer input
def get_positive_integer_input(prompt):
    while True:
        try:
            value = int(input(prompt))
            if value <= 0 :
                raise ValueError
            return value
        except ValueError:
            print(Fore.RED + Style.BRIGHT + 'Invalid input. Please enter a positive integer.')


# Function to generate the presentation
def generate_presentation():
    prs = create_presentation()
    while True:
        try:
            print()
            number_of_slides = get_positive_integer_input(Fore.YELLOW + Style.BRIGHT + 'Enter the number of slides: ')

            for i in range(1, number_of_slides + 1):
                print()
                print(Fore.LIGHTYELLOW_EX + Style.BRIGHT + f"Slide - {i}")
                type_of_slide = get_num_between_1_and_3(
                    Fore.MAGENTA + Style.BRIGHT + f'''Enter the layout of slide {i}:
1. Title slide
2. Graph or Table slide
3. Image slide 
                    option - ''')
                if type_of_slide not in [1, 2, 3]:
                    print(Fore.RED + Style.BRIGHT + "Invalid slide layout choice. Please try again.")
                    return

                slide_layouts = prs.slide_layouts
                slide_layout = slide_layouts[type_of_slide - 1]
                slide = add_slide(prs, slide_layout)

                if type_of_slide == 1:
                    if i == number_of_slides - 1:
                        add_title_slide(slide, "Thank You")
                    else:
                        slide_title = input(Fore.YELLOW + Style.BRIGHT + 'Input title: ')
                        add_title_slide(slide, slide_title)

                elif type_of_slide == 2:
                    add_title_content_slide(slide)

                elif type_of_slide == 3:
                    add_section_header_slide(slide)

                # Set slide color
                print()
                while True:
                    slide_color = input(Fore.YELLOW + Style.BRIGHT + '''Enter the color for the slide 
(color name or hexadecimal code): ''')

                    if slide_color.startswith('#'):
                        slide_color = slide_color[1:]  # Remove the '#' prefix if present

                    try:
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor.from_string(slide_color)

                        # Set text color based on background color
                        text_color = determine_text_color(webcolors.hex_to_rgb(slide_color))
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.color.rgb = RGBColor.from_string(text_color)

                        # Set legend text color for graphs
                        for shape in slide.shapes:
                            if shape.has_chart:
                                chart = shape.chart
                                if chart.has_legend:
                                    legend = chart.legend
                                    legend_text_color = determine_text_color(webcolors.hex_to_rgb(slide_color))
                                    legend.font.color.rgb = RGBColor.from_string(legend_text_color)

                        break
                    except ValueError:
                        # Handle color names
                        try:
                            rgb_color = webcolors.name_to_rgb(slide_color)
                            slide.background.fill.solid()
                            slide.background.fill.fore_color.rgb = RGBColor(rgb_color.red, rgb_color.green,
                                                                            rgb_color.blue)

                            # Set text color based on background color
                            text_color = determine_text_color((rgb_color.red, rgb_color.green, rgb_color.blue))
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for paragraph in shape.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.color.rgb = RGBColor.from_string(text_color)

                            # Set legend text color for graphs
                            for shape in slide.shapes:
                                if shape.has_chart:
                                    chart = shape.chart
                                    if chart.has_legend:
                                        legend = chart.legend
                                        legend_text_color = determine_text_color(
                                            (rgb_color.red, rgb_color.green, rgb_color.blue))
                                        legend.font.color.rgb = RGBColor.from_string(legend_text_color)

                            break
                        except ValueError:
                            print(Fore.RED + Style.BRIGHT + "Invalid color entered.")

            # Add a "Thank You" slide at the end
            thank_you_slide_layout = prs.slide_layouts[0]  # Choose a suitable layout for the "Thank You" slide
            thank_you_slide = add_slide(prs, thank_you_slide_layout)
            add_title_slide(thank_you_slide, "Thank You")

        except ValueError:
            print(Fore.RED + Style.BRIGHT + "Invalid input. Please enter a valid number of slides.")
            return
        break

    # Save the presentation
    file_name = f"sales_report.pptx"
    prs.save(f"") # file path in this format(keep the variable) C:\\Users\\OneDrive\\Desktop\\code\\{file_name}
    os.startfile(f"")  # file path in this format(keep the variable) C:\\Users\\OneDrive\\Desktop\\code\\{file_name}

print(Fore.LIGHTBLUE_EX + Style.BRIGHT + '''Hey want to create a PPT?
Dont worry we got you!!!
Just type the required fields below.''')
generate_presentation()
